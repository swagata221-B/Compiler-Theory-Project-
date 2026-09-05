#!/usr/bin/env python3
"""Build Presentation 2 slides from the MiniPascal compiler project."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("MiniPascal_CSE303_Presentation2.pptx")

PAPER = RGBColor(0xF3, 0xEE, 0xE4)
INK = RGBColor(0x1C, 0x19, 0x14)
SOFT = RGBColor(0x4A, 0x45, 0x3C)
RULE = RGBColor(0xCB, 0xBF, 0xAA)
BURG = RGBColor(0x8B, 0x29, 0x42)
GOLD = RGBColor(0xB0, 0x89, 0x48)
CHIP = RGBColor(0xEF, 0xE6, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
W = Inches(13.333)
H = Inches(7.5)


def set_run(run, text, *, size=18, bold=False, color=INK, font="Calibri", italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_text_box(slide, l, t, w, h, text, *, size=18, bold=False, color=INK, font="Calibri", italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        set_run(run, line, size=size, bold=bold, color=color, font=font, italic=italic)
    return box


def fill_slide(slide, color=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    # send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return bg


def bar(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(slide)
    bar(slide, 0, 0, Inches(0.12), H, BURG)
    add_text_box(
        slide,
        Inches(0.45),
        Inches(0.18),
        Inches(10),
        Inches(0.32),
        "CSE303  ·  MINIPASCAL COMPILER  ·  PRESENTATION 2",
        size=11,
        color=BURG,
        font="Calibri",
        bold=True,
    )
    return slide


def heading(slide, title, subtitle=None):
    add_text_box(slide, Inches(0.45), Inches(0.48), Inches(12.3), Inches(0.7), title, size=32, bold=True, font="Georgia", color=INK)
    if subtitle:
        add_text_box(slide, Inches(0.45), Inches(1.15), Inches(12.3), Inches(0.4), subtitle, size=16, color=SOFT, italic=True)


def bullets(slide, items, l=Inches(0.5), t=Inches(1.75), w=Inches(12.2), size=18, color=INK):
    box = slide.shapes.add_textbox(l, t, w, Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = item[0] if isinstance(item, tuple) else 0
        text = item[1] if isinstance(item, tuple) else item
        p.space_after = Pt(8)
        run = p.add_run()
        set_run(run, text, size=size - (2 if p.level else 0), color=color, font="Calibri")
    return box


def card(slide, l, t, w, h, title, body, accent=BURG):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = CHIP
    sh.line.color.rgb = RULE
    bar(slide, l, t, Inches(0.08), h, accent)
    add_text_box(slide, l + Inches(0.22), t + Inches(0.16), w - Inches(0.36), Inches(0.4), title, size=15, bold=True, font="Georgia", color=accent)
    add_text_box(slide, l + Inches(0.22), t + Inches(0.58), w - Inches(0.36), h - Inches(0.72), body, size=14, color=SOFT)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s)
    bar(s, 0, 0, Inches(0.22), H, BURG)
    bar(s, 0, H - Inches(0.18), W, Inches(0.18), GOLD)
    add_text_box(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4), "CSE303  ·  COMPILER DESIGN  ·  PRESENTATION 2", size=14, bold=True, color=BURG)
    add_text_box(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.8), "MiniPascal frontend", size=48, bold=True, font="Georgia")
    add_text_box(s, Inches(0.7), Inches(3.7), Inches(11), Inches(1.0), "A Flex scanner, a Bison parser, and an AST in C —\nbuilt first as a compiler, then described by these slides.", size=22, color=SOFT)
    add_text_box(s, Inches(0.7), Inches(6.4), Inches(11), Inches(0.4), "Scope of this talk: timeline 2 only. Semantics and code generation are later.", size=14, italic=True, color=SOFT)
    notes(s, "Open by saying the compiler exists as a project. These slides report on that project; they are not the project.")

    # 2 agenda
    s = new_slide(prs)
    heading(s, "Agenda", "What we will show from the repository")
    bullets(s, [
        "Where timeline 2 sits in the compiler",
        "What the project actually implements — and what it refuses to do yet",
        "Lexer: Flex patterns, comments, the 1..10 trap",
        "Parser: Bison grammar, precedence, AST actions",
        "Live proof: samples/gcd.pas and samples/broken.pas",
        "Handoff to timeline 3",
    ])
    notes(s, "Keep this under a minute. The demo is the close.")

    # 3 timeline
    s = new_slide(prs)
    heading(s, "Timeline", "The project stops after phase 2")
    phases = [
        ("01", "Language and architecture", "Context", "Named the MiniPascal subset and the four-phase pipeline.", RULE),
        ("02", "Lexer and parser", "This talk", "Flex scanner, Bison parser, AST, diagnostics with line and column.", BURG),
        ("03", "Semantic analysis", "Later", "Symbol tables, types, declaration-before-use. Input is today's AST.", GOLD),
        ("04", "Code generation", "Later", "Walk a typed tree and emit a simple target. Not in this repository.", SOFT),
    ]
    for i, (num, title, tag, body, accent) in enumerate(phases):
        x = Inches(0.45) + i * Inches(3.2)
        sh = slide_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.85), Inches(3.0), Inches(4.3))
        slide_card.adjustments[0] = 0.05
        slide_card.fill.solid()
        slide_card.fill.fore_color.rgb = CHIP if i == 1 else PAPER
        slide_card.line.color.rgb = accent
        add_text_box(s, x + Inches(0.18), Inches(2.05), Inches(2.6), Inches(0.4), num, size=20, bold=True, font="Georgia", color=accent)
        add_text_box(s, x + Inches(0.18), Inches(2.5), Inches(2.6), Inches(0.35), tag.upper(), size=11, bold=True, color=accent)
        add_text_box(s, x + Inches(0.18), Inches(2.95), Inches(2.6), Inches(1.1), title, size=18, bold=True, font="Georgia")
        add_text_box(s, x + Inches(0.18), Inches(4.15), Inches(2.6), Inches(1.7), body, size=13, color=SOFT)
    notes(s, "Point at column 02. Say P3 and P4 are empty folders in the plan, not missing features of this demo.")

    # 4 in / out of scope
    s = new_slide(prs)
    heading(s, "Scope of the project", "Built for timeline 2 — nothing more")
    card(s, Inches(0.45), Inches(1.8), Inches(6.0), Inches(4.8), "In the repository",
         "Flex scanner with line and column.\n\nBison LALR parser for MiniPascal.\n\nAST nodes built in the grammar actions.\n\nLex and parse errors with a source position.\n\nCLI: ./minipascal and ./minipascal --tokens")
    card(s, Inches(6.75), Inches(1.8), Inches(6.0), Inches(4.8), "Not in this phase",
         "No symbol table.\n\nNo type checking.\n\nNo “undeclared identifier” errors.\n\nNo intermediate code.\n\nNo target machine or interpreter.\n\nUsing an undeclared name is still a legal syntax tree.", accent=GOLD)
    notes(s, "If asked why x := y parses with no declaration of y: that is a semantic question. The parser's job is structure.")

    # 5 pipeline
    s = new_slide(prs)
    heading(s, "Pipeline we implemented", "flex lexer.l   +   bison parser.y   +   gcc")
    steps = [
        ("Source", ".pas text"),
        ("Lexer", "lex.yy.c"),
        ("Tokens", "kind + span"),
        ("Parser", "parser.tab.c"),
        ("AST", "Program node"),
    ]
    for i, (title, sub) in enumerate(steps):
        x = Inches(0.5) + i * Inches(2.5)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.4), Inches(2.15), Inches(1.7))
        box.adjustments[0] = 0.12
        box.fill.solid()
        box.fill.fore_color.rgb = BURG if i in (1, 3) else CHIP
        box.line.fill.background()
        fg = WHITE if i in (1, 3) else INK
        add_text_box(s, x + Inches(0.12), Inches(2.7), Inches(1.9), Inches(0.55), title, size=20, bold=True, font="Georgia", color=fg, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.12), Inches(3.35), Inches(1.9), Inches(0.45), sub, size=13, color=fg, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text_box(s, x + Inches(2.0), Inches(2.9), Inches(0.45), Inches(0.5), "→", size=22, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.5), Inches(12), Inches(1.8),
                 "Flex generates the scanner. Bison generates the LALR parser.\nWe write the patterns, the grammar, and the AST actions.\nErrors print a line and column. We do not type-check in this phase.",
                 size=16, color=SOFT)
    notes(s, "Draw the arrow with your hand: characters, words, tree. Stop before 'meaning' and 'code'.")

    # 6 layout
    s = new_slide(prs)
    heading(s, "What is in the project", "The slides sit beside the compiler, not around it")
    bullets(s, [
        "compiler/lexer.l — Flex scanner",
        "compiler/parser.y — Bison grammar and AST actions",
        "compiler/ast.c — Program, statements, expressions",
        "./minipascal file.pas   and   ./minipascal --tokens file.pas",
        "samples/gcd.pas, factorial.pas, bubble.pas, broken.pas",
        "presentation/MiniPascal_CSE303_Presentation2.pptx — this deck",
    ], t=Inches(1.7))
    add_text_box(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
                 "make test checks the sample programs and that 1..10 tokenizes as a range.",
                 size=15, italic=True, color=SOFT)
    notes(s, "Open the tree on the laptop if you can. Point at presentation/ last, so it is clear the PPTX is an output.")

    # 7 lexer
    s = new_slide(prs)
    heading(s, "Lexer", "compiler/lexer.l · Flex")
    card(s, Inches(0.45), Inches(1.8), Inches(4.0), Inches(4.7), "Recognizes",
         "Keywords, case-insensitive.\n\nIdentifiers and literals.\n\n:=  <>  <=  >=  ..\n\n{ }  (* *)  and line comments.\n\nPascal strings: 'it''s'.")
    card(s, Inches(4.65), Inches(1.8), Inches(4.0), Inches(4.7), "Records",
         "Token kind.\n\nOriginal lexeme.\n\nStart line and column.\n\nEnd line and column.\n\nNumeric value for integers and reals.", accent=GOLD)
    card(s, Inches(8.85), Inches(1.8), Inches(4.0), Inches(4.7), "Rejects",
         "Unexpected characters.\n\nUnterminated strings.\n\nUnterminated comments.\n\nA real with a broken exponent.\n\nEach error keeps a source position.", accent=SOFT)
    notes(s, "The one trick to mention: after digits, '..' must not become a real number. 1..10 is three tokens.")

    # 8 tokens from gcd
    s = new_slide(prs)
    heading(s, "Tokens from samples/gcd.pas", "First lines of a real program in the repo")
    rows = [
        ("Line", "Kind", "Lexeme"),
        ("1", "PROGRAM", "program"),
        ("1", "IDENT", "Gcd"),
        ("2", "VAR", "var"),
        ("3", "IDENT", "a , b , t"),
        ("6", "WHILE", "while"),
        ("6", "NEQ", "<>"),
        ("9", "MOD", "mod"),
        ("12", "STRING_LIT", "gcd = "),
        ("13", "DOT", "."),
    ]
    top = Inches(1.7)
    widths = [Inches(1.6), Inches(3.4), Inches(5.2)]
    lefts = [Inches(0.5), Inches(2.2), Inches(5.8)]
    for r, row in enumerate(rows):
        y = top + r * Inches(0.42)
        bgc = BURG if r == 0 else (CHIP if r % 2 == 0 else PAPER)
        fg = WHITE if r == 0 else INK
        bar(s, Inches(0.45), y, Inches(11.8), Inches(0.42), bgc)
        for c, cell in enumerate(row):
            add_text_box(s, lefts[c], y + Inches(0.05), widths[c], Inches(0.32), cell, size=14, bold=(r == 0), color=fg, font="Consolas")
    notes(s, "Run ./minipascal --tokens samples/gcd.pas for a live dump.")

    # 9 parser
    s = new_slide(prs)
    heading(s, "Parser", "compiler/parser.y · Bison LALR(1)")
    bullets(s, [
        "Each nonterminal is a Bison rule: program, declarations, statement, expression…",
        "After an identifier, ':=' is assignment and '(' is a call",
        "else binds to the nearest if (%nonassoc THEN / ELSE)",
        "Precedence: or, and, relops, + −, * / div mod, then unary not / + / −",
        "yyerror prints the line and column; error_count becomes the exit status",
        "Actions allocate AST nodes. Later phases walk that tree.",
    ], t=Inches(1.7))
    notes(s, "You can show parseAssignOrCall if someone asks how writeln(a) is not an assignment.")

    # 10 grammar
    s = new_slide(prs)
    heading(s, "Grammar we accept", "Teaching MiniPascal, not full Pascal")
    add_text_box(
        s,
        Inches(0.5),
        Inches(1.7),
        Inches(12.3),
        Inches(5.2),
        "program      ::=  program ident ';' declarations subprograms compound '.'\n"
        "declarations ::=  [ var (ident-list ':' type ';')+ ]\n"
        "type         ::=  integer | real | boolean | char\n"
        "             |    array '[' int '..' int ']' of standard-type\n"
        "subprogram   ::=  function ident args ':' type ';' … compound\n"
        "             |    procedure ident args ';' … compound\n"
        "statement    ::=  assign | call | begin-end | if | while | for | return\n"
        "if           ::=  if expr then stmt [ else stmt ]",
        size=16,
        font="Consolas",
        color=INK,
    )
    notes(s, "This is the language the parser implements. Full Pascal (records, pointers, units) is out of scope for the course dialect.")

    # 11 AST
    s = new_slide(prs)
    heading(s, "AST we hand to timeline 3", "From samples/gcd.pas after a successful parse")
    add_text_box(
        s,
        Inches(0.5),
        Inches(1.7),
        Inches(12.3),
        Inches(5.2),
        "Program Gcd\n"
        "  Declarations\n"
        "    VarDecl   a, b, t : integer\n"
        "  BeginEnd\n"
        "    Call      read(a, b)\n"
        "    While     b <> 0\n"
        "      Assign  t := b\n"
        "      Assign  b := a mod b\n"
        "      Assign  a := t\n"
        "    Call      writeln('gcd = ', a)",
        size=18,
        font="Consolas",
        color=INK,
    )
    notes(s, "This tree is the only artifact presentation 3 should need. No need to re-lex.")

    # 12 demo happy
    s = new_slide(prs)
    heading(s, "Demo — accepted program", "./minipascal samples/gcd.pas")
    add_text_box(
        s,
        Inches(0.5),
        Inches(1.7),
        Inches(6.2),
        Inches(5.0),
        "program Gcd;\n"
        "var\n"
        "  a, b, t: integer;\n"
        "begin\n"
        "  read(a, b);\n"
        "  while b <> 0 do\n"
        "  begin\n"
        "    t := b;\n"
        "    b := a mod b;\n"
        "    a := t\n"
        "  end;\n"
        "  writeln('gcd = ', a)\n"
        "end.",
        size=16,
        font="Consolas",
    )
    card(s, Inches(7.1), Inches(1.8), Inches(5.6), Inches(4.6), "What the frontend reports",
         "Exit code 0.\n\nToken stream for the whole file.\n\nA Program node named Gcd.\n\nA While whose test is the <> operator.\n\nmod is a keyword token, then a binary operator in the tree.\n\nNo diagnostics.")
    notes(s, "Run this live. Also mention factorial.pas (function + locals) if there is time.")

    # 13 demo broken
    s = new_slide(prs)
    heading(s, "Demo — rejected program", "./minipascal samples/broken.pas")
    add_text_box(
        s,
        Inches(0.5),
        Inches(1.7),
        Inches(6.2),
        Inches(3.4),
        "program Broken\n"
        "var\n"
        "  x: integer;\n"
        "begin\n"
        "  x := ;\n"
        "  writeln(x)\n"
        "end",
        size=18,
        font="Consolas",
    )
    card(s, Inches(7.1), Inches(1.8), Inches(5.6), Inches(4.6), "Diagnostics the parser emits",
         "2:1  Expected ';' after the program name\n\n5:8  Expected an expression\n\n8:1  A MiniPascal program must end with '.'\n\nExit code 1.\n\nA partial tree may still exist after recovery.")
    notes(s, "This is the close of the technical demo. Three errors, three positions, no crash.")

    # 14 run
    s = new_slide(prs)
    heading(s, "How we run the project", "VS Code terminal, or any Unix shell / WSL")
    bullets(s, [
        "sudo apt install build-essential flex bison",
        "make",
        "./minipascal samples/gcd.pas",
        "./minipascal --tokens samples/gcd.pas",
        "./minipascal samples/broken.pas",
        "make test",
    ], t=Inches(1.75), size=20)
    notes(s, "Run these commands in VS Code's terminal. No Node.js.")

    # 15 next
    s = new_slide(prs)
    heading(s, "What timeline 3 will take from this", "The AST is the hand-off")
    bullets(s, [
        "Walk Program, Subprogram, and Compound with a scoped symbol table",
        "Insert every VarDecl and parameter before checking uses",
        "Type-check Assign, Call, If, While, For, and operators",
        "Reject undeclared names and mismatched operands",
        "Leave code generation for timeline 4 — still a tree walk, now on a typed AST",
    ], t=Inches(1.75))
    add_text_box(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
                 "We will not rewrite the lexer or the parser for that talk.",
                 size=16, italic=True, color=BURG)
    notes(s, "Close the technical story here. Invite questions on the frontend only.")

    # 16 end
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide(s)
    bar(s, 0, 0, Inches(0.22), H, BURG)
    bar(s, 0, H - Inches(0.18), W, Inches(0.18), GOLD)
    add_text_box(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.2), "Questions", size=48, bold=True, font="Georgia")
    add_text_box(s, Inches(0.7), Inches(3.6), Inches(11.5), Inches(1.2),
                 "The compiler is the project.\nThese slides only explain timeline 2 of that compiler.",
                 size=22, color=SOFT)
    add_text_box(s, Inches(0.7), Inches(6.3), Inches(11), Inches(0.4),
                 "CSE303 MiniPascal  ·  Presentation 2  ·  lexer + parser + AST",
                 size=14, color=BURG, bold=True)
    notes(s, "If someone asks about codegen, say it is timeline 4 and the AST is already the input.")

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
